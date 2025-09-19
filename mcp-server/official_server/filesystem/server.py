"""
Filesystem FastMCP Server (sandboxed) - 标准化版本

功能：
- 列目录、查询文件信息、读写文本、读写二进制（以 base64 返回）、创建目录、复制/移动/删除文件
- 提取 Office 文本：PDF / DOCX / PPTX
- 使用标准化MCP响应格式

依赖（按需可选安装）：
- PDF:  "pypdf"
- DOCX: "python-docx"  
- PPTX: "python-pptx"

运行（开发调试，自动装依赖）：
    uv run mcp dev examples/servers/filesystem/server.py \
        --with pypdf python-docx python-pptx

生产建议：使用 streamable-http
    uv run examples/servers/filesystem/server.py

环境变量：
- FILESYSTEM_BASE_DIR：沙箱根目录（默认：用户家目录）
"""

from __future__ import annotations

import base64
import json
import os
import shutil
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional

from pydantic import BaseModel, Field
from mcp.server.fastmcp import FastMCP

# 导入标准化组件
import sys
sys.path.append(os.path.join(os.path.dirname(__file__), '../..'))
from core import (
    MCPResponseBuilder, OperationType, ResponseStatus, create_file_info,
    FileSystemTool, quick_success, quick_error
)

# -----------------------------------------------------------------------------
# 配置与沙箱  
# -----------------------------------------------------------------------------

def load_resource_config() -> dict:
    """加载资源配置文件"""
    config_file = "resource_config.json"
    default_config = {
        "base_directory": "./mcp_resources",
        "auto_create_dirs": True,
        "max_file_size_mb": 100
    }
    
    if os.path.exists(config_file):
        try:
            with open(config_file, 'r', encoding='utf-8') as f:
                config = json.load(f)
                default_config.update(config)
        except Exception:
            pass
    
    return default_config

def get_base_dir() -> Path:
    """获取基础目录，优先使用MCP_BASE_DIR环境变量，然后配置文件，最后默认值"""
    # 1. 优先使用MCP_BASE_DIR环境变量（由launcher设置）
    mcp_base = os.environ.get("MCP_BASE_DIR")
    if mcp_base:
        try:
            p = Path(mcp_base).expanduser().resolve()
            if not p.exists():
                p.mkdir(parents=True, exist_ok=True)
            return p
        except Exception:
            pass
    
    # 2. 尝试从配置文件获取
    try:
        config = load_resource_config()
        config_path = config.get("base_directory")
        if config_path:
            p = Path(config_path).expanduser().resolve()
            if not p.exists():
                p.mkdir(parents=True, exist_ok=True)
            return p
    except Exception:
        pass
    
    # 3. 尝试从FILESYSTEM_BASE_DIR环境变量获取
    env_path = os.environ.get("FILESYSTEM_BASE_DIR")
    if env_path:
        try:
            p = Path(env_path).expanduser().resolve()
            if not p.exists():
                p.mkdir(parents=True, exist_ok=True)
            return p
        except Exception:
            pass
    
    # 4. 默认使用当前目录下的mcp_workspace
    default_path = Path(os.getcwd()) / "mcp_workspace"
    default_path.mkdir(parents=True, exist_ok=True)
    return default_path.resolve()


BASE_DIR: Path = get_base_dir()
RESOURCE_CONFIG = load_resource_config()

def resolve_in_sandbox(user_path: str, session_id: str = None) -> Path:
    """将用户提供路径解析到沙箱内，拒绝越权路径。

    - 支持绝对/相对路径；相对路径相对于BASE_DIR
    - 解析后必须保证在BASE_DIR沙箱目录内
    """
    # 直接使用BASE_DIR作为沙箱根目录
    sandbox_root = BASE_DIR
    
    candidate = Path(user_path)
    if not candidate.is_absolute():
        candidate = sandbox_root / candidate
    abs_path = candidate.expanduser().resolve()
    
    try:
        # 允许目标等于sandbox_root本身或其子路径
        if abs_path == sandbox_root or sandbox_root in abs_path.parents:
            return abs_path
    except Exception:
        pass
    raise ValueError(f"Path out of sandbox: {abs_path} (sandbox: {sandbox_root})")

# -----------------------------------------------------------------------------
# 数据模型（结构化输出）
# -----------------------------------------------------------------------------

class FileInfo(BaseModel):
    path: str = Field(description="绝对路径")
    relative: str = Field(description="相对 BASE_DIR 的路径")
    name: str
    is_dir: bool
    size: int | None = None
    mtime: str | None = None


class ListDirRequest(BaseModel):
    path: str = Field(description="目录路径（可相对沙箱）")
    recursive: bool = Field(default=False)
    pattern: str | None = Field(default=None, description="glob 过滤，如 *.txt")
    files_only: bool = Field(default=True)


class ListDirResult(BaseModel):
    base_dir: str
    entries: list[FileInfo]


class ReadTextRequest(BaseModel):
    path: str
    encoding: str = Field(default="utf-8")


class WriteTextRequest(BaseModel):
    path: str
    content: str = Field(default="", description="文件内容，默认为空字符串")
    encoding: str = Field(default="utf-8")
    append: bool = Field(default=False)


class ReadBinaryRequest(BaseModel):
    path: str
    max_bytes: int | None = Field(default=None, description="限制最大读取字节，None 表示不限")


class ReadBinaryResult(BaseModel):
    base64: str
    size: int
    mime_type: str = Field(default="application/octet-stream")


class WriteBinaryRequest(BaseModel):
    path: str
    base64: str
    overwrite: bool = Field(default=True)


class MoveCopyRequest(BaseModel):
    src: str
    dst: str
    overwrite: bool = Field(default=False)


class MkdirRequest(BaseModel):
    path: str
    parents: bool = Field(default=True)
    exist_ok: bool = Field(default=True)


class DeleteRequest(BaseModel):
    path: str
    recursive: bool = Field(default=False)


class OfficeReadRequest(BaseModel):
    path: str


def to_file_info(p: Path, sandbox_root: Path = None) -> FileInfo:
    stat = p.stat()
    if sandbox_root is None:
        sandbox_root = BASE_DIR
    
    # 计算相对路径
    try:
        if p == sandbox_root:
            rel = Path(".")
        else:
            rel = p.relative_to(sandbox_root)
    except ValueError:
        # 如果路径不在sandbox_root下，尝试相对于BASE_DIR
        try:
            if p == BASE_DIR:
                rel = Path(".")
            else:
                rel = p.relative_to(BASE_DIR)
        except ValueError:
            rel = Path(".")
    
    return FileInfo(
        path=str(p),
        relative=str(rel),
        name=p.name,
        is_dir=p.is_dir(),
        size=None if p.is_dir() else int(stat.st_size),
        mtime=datetime.fromtimestamp(stat.st_mtime).isoformat(),
    )

# -----------------------------------------------------------------------------
# 服务器与工具
# -----------------------------------------------------------------------------

# 创建MCP服务器和工具实例
mcp = FastMCP("filesystem")

# 创建文件系统工具实例
fs_tool = FileSystemTool("filesystem", BASE_DIR)

@mcp.tool()
def get_base() -> Dict[str, Any]:
    """获取当前文件系统沙箱的根目录路径。
    
    返回MCP服务器能够访问的基础目录路径，所有文件操作都限制在此目录下。
    这对于了解可操作的文件范围很有用。
    
    Returns:
        包含基础目录信息的响应，其中data字段包含base_directory路径
    """
    return quick_success(
        tool_name="get_base",
        operation=OperationType.QUERY,
        message="成功获取基础目录",
        data={"base_directory": str(BASE_DIR)}
    )

# 专用于直接调用的路径列表函数 - 不向AI提供
# 注意：这个函数故意不使用 @mcp.tool() 装饰器，只能通过直接HTTP调用
def list_all_paths(session_id: str = None) -> Dict[str, Any]:
    """获取base_dir下所有文件和文件夹的绝对路径列表。
    
    这是一个专门用于直接工具调用的函数，不会被AI智能任务执行器调用。
    返回base_dir及其子目录下所有文件和文件夹的绝对路径。
    路径格式会根据操作系统自动调整（Windows/Linux）。
    
    重要：此函数故意不注册为MCP工具，只能通过客户端直接调用接口使用。
    """
    builder = MCPResponseBuilder("list_all_paths")
    
    try:
        paths = []
        root = BASE_DIR
        
        # 添加根目录本身
        paths.append(str(root))
        
        # 遍历所有子项（文件和文件夹）
        for path in root.rglob("*"):
            # 忽略.useit文件夹及其内容
            if ".useit" not in path.parts:
                paths.append(str(path))
        
        # 排序：目录在前，文件在后，同类型按名称排序
        def sort_key(p: str) -> tuple[bool, str]:
            path_obj = Path(p)
            return (not path_obj.is_dir(), p.lower())
        
        paths.sort(key=sort_key)
        
        return builder.success(
            operation=OperationType.QUERY,
            message=f"成功列出 {len(paths)} 个路径",
            data={"paths": paths, "total": len(paths)}
        ).to_dict()
    
    except Exception as e:
        return builder.error(
            operation=OperationType.QUERY,
            message="列出路径失败",
            error_details=str(e)
        ).to_dict()

@mcp.tool()
def list_dir(
    path: str = Field(".", description="要列出的目录路径。使用'.'表示当前工作目录，可以是相对路径或绝对路径"),
    recursive: bool = Field(False, description="是否递归搜索子目录。True表示会列出所有子目录中的文件"),
    pattern: Optional[str] = Field(None, description="文件名过滤模式，支持glob语法。例如：'*.txt'只显示文本文件，'test*'显示以test开头的文件"),
    files_only: bool = Field(True, description="是否只显示文件。True只显示文件，False同时显示文件和目录")
) -> Dict[str, Any]:
    """列出指定目录中的文件和文件夹。
    
    这是查看目录内容的主要工具。支持递归搜索、文件名过滤等高级功能。
    默认情况下只显示文件，如果需要查看目录结构请设置files_only=False。
    
    Args:
        path: 要列出的目录路径。使用'.'表示当前工作目录
        recursive: 是否递归搜索子目录中的所有文件
        pattern: 文件名过滤模式，支持通配符语法
        files_only: 是否只显示文件（不显示目录）
    
    Returns:
        包含目录内容列表的响应，每个文件/目录都有完整的元数据信息
    
    Examples:
        - list_dir() 或 list_dir(".") : 列出当前目录的所有文件
        - list_dir("docs", recursive=True) : 递归列出docs目录下的所有文件
        - list_dir(pattern="*.py") : 只显示Python文件
        - list_dir(files_only=False) : 显示文件和目录
    """
    builder = MCPResponseBuilder("list_dir")
    
    try:
        # 创建请求对象
        req = ListDirRequest(
            path=path,
            recursive=recursive,
            pattern=pattern,
            files_only=files_only
        )
        
        # 如果请求路径为"."或空，使用BASE_DIR并启用递归搜索
        if req.path in (".", "", "/"):
            root = BASE_DIR
            use_recursive = True
            max_items = 300
        else:
            root = resolve_in_sandbox(req.path)
            use_recursive = req.recursive
            max_items = None
        
        if not root.exists():
            return builder.error(
                operation=OperationType.QUERY,
                message="目录不存在",
                error_details=f"目录 '{req.path}' 不存在"
            ).to_dict()
            
        if not root.is_dir():
            return builder.error(
                operation=OperationType.QUERY,
                message="不是目录",
                error_details=f"'{req.path}' 不是一个目录"
            ).to_dict()

        def iter_paths() -> list[Path]:
            paths = []
            count = 0
            
            if req.pattern:
                if use_recursive:
                    pattern_paths = root.rglob(req.pattern)
                else:
                    pattern_paths = root.glob(req.pattern)
            else:
                if use_recursive:
                    pattern_paths = root.rglob("*")
                else:
                    pattern_paths = root.glob("*")
            
            for p in pattern_paths:
                # 忽略.useit文件夹及其内容
                if ".useit" in p.parts:
                    continue
                    
                paths.append(p)
                count += 1
                
                # 如果设置了最大数量限制，则停止收集
                if max_items and count >= max_items:
                    break
            
            return paths

        paths = iter_paths()
        if req.files_only:
            paths = [p for p in paths if p.is_file()]

        # 使用BASE_DIR作为sandbox_root
        entries = [to_file_info(p, BASE_DIR) for p in paths]
        
        # 按类型和名称排序：先目录后文件，同类型按名称排序
        entries.sort(key=lambda x: (not x.is_dir, x.name.lower()))
        
        result_data = ListDirResult(base_dir=str(BASE_DIR), entries=entries)
        
        # 统计信息
        file_count = sum(1 for e in entries if not e.is_dir)
        dir_count = sum(1 for e in entries if e.is_dir)
        total_size = sum(e.size or 0 for e in entries if not e.is_dir)
        
        return builder.success(
            operation=OperationType.QUERY,
            message=f"成功列出目录 '{req.path}'，共 {len(entries)} 个条目（{file_count} 个文件，{dir_count} 个目录）",
            data={
                **result_data.model_dump(),
                "summary": {
                    "total_entries": len(entries),
                    "file_count": file_count,
                    "directory_count": dir_count,
                    "total_size": total_size,
                    "search_params": {
                        "recursive": use_recursive,
                        "pattern": req.pattern,
                        "files_only": req.files_only
                    }
                }
            }
        ).to_dict()
        
    except Exception as e:
        return builder.error(
            operation=OperationType.QUERY,
            message="列目录失败",
            error_details=str(e)
        ).to_dict()



@mcp.tool()
def read_text(
    path: str = Field(..., description="要读取的文本文件路径，可以是相对路径或绝对路径。例如：'document.txt', 'folder/file.md'"),
    encoding: str = Field("utf-8", description="文件的字符编码格式。常用值：'utf-8'(默认), 'gbk', 'ascii', 'latin-1'")
) -> Dict[str, Any]:
    """读取文本文件的完整内容。
    
    用于读取各种文本格式文件，包括源代码、配置文件、文档等。
    支持不同的字符编码，默认使用UTF-8编码。
    
    Args:
        path: 要读取的文本文件路径
        encoding: 文件的字符编码格式，确保正确解码文件内容
    
    Returns:
        包含文件完整内容的响应，同时提供文件大小、行数等元数据
    
    Examples:
        - read_text("README.md") : 读取README文件
        - read_text("config.json") : 读取JSON配置文件
        - read_text("chinese.txt", encoding="gbk") : 读取GBK编码的中文文件
    """
    builder = MCPResponseBuilder("read_text")
    
    try:
        # 创建请求对象
        req = ReadTextRequest(
            path=path,
            encoding=encoding
        )
        
        p = resolve_in_sandbox(req.path)
        if not p.exists() or not p.is_file():
            return builder.error(
                operation=OperationType.READ,
                message="文件不存在或不是文件",
                error_details=f"路径 '{req.path}' 不是一个有效文件"
            ).to_dict()
            
        content = p.read_text(encoding=req.encoding)
        
        content_bytes = content.encode(req.encoding)
        return builder.success(
            operation=OperationType.READ,
            message=f"成功读取文件 '{req.path}'",
            data={
                "content": content,
                "size": len(content_bytes),
                "encoding": req.encoding,
                "line_count": len(content.splitlines()),
                "character_count": len(content)
            }
        ).to_dict()
        
    except UnicodeDecodeError as e:
        return builder.error(
            operation=OperationType.READ,
            message="文件编码错误",
            error_details=f"无法使用 {req.encoding} 编码读取文件: {str(e)}"
        ).to_dict()
        
    except Exception as e:
        return builder.error(
            operation=OperationType.READ,
            message="读取文件失败",
            error_details=str(e)
        ).to_dict()

@mcp.tool()
def write_file(
    path: str = Field(..., description="要写入的文件路径。如果目录不存在会自动创建。例如：'output.txt', 'data/result.json'"),
    content: str = Field(..., description="要写入的文本内容，不能为空"),
    encoding: str = Field("utf-8", description="文件的字符编码格式。常用值：'utf-8'(默认), 'gbk', 'ascii'"),
    append: bool = Field(False, description="写入模式。False表示覆盖写入(默认)，True表示追加到文件末尾")
) -> Dict[str, Any]:
    """创建新文件或写入/追加内容到文本文件。
    
    这是文件写入的主要工具。可以创建新文件、覆盖现有文件内容或追加内容到文件末尾。
    如果目标目录不存在，会自动创建所需的目录结构。
    
    Args:
        path: 要写入的文件路径，支持相对和绝对路径
        content: 要写入的文本内容，可以是任何文本格式
        encoding: 文件的字符编码格式
        append: 写入模式选择（覆盖或追加）
    
    Returns:
        包含操作结果的响应，显示文件是否为新创建、文件大小等信息
    
    Examples:
        - write_file("hello.txt", "Hello World") : 创建文件并写入内容
        - write_file("log.txt", "新日志\\n", append=True) : 追加日志到现有文件
        - write_file("empty.txt") : 创建空文件
        - write_file("data.json", json_content) : 保存JSON数据
    """
    builder = MCPResponseBuilder("write_file")
    
    try:
        # 创建请求对象
        req = WriteTextRequest(
            path=path,
            content=content,
            encoding=encoding,
            append=append
        )
        
        p = resolve_in_sandbox(req.path)
        p.parent.mkdir(parents=True, exist_ok=True)
        
        is_new_file = not p.exists()
        if req.append and p.exists():
            with p.open("a", encoding=req.encoding) as f:
                f.write(req.content)
            operation_type = OperationType.UPDATE
            message = f"成功追加内容到文件 '{req.path}'"
        else:
            with p.open("w", encoding=req.encoding) as f:
                f.write(req.content)
            operation_type = OperationType.CREATE if is_new_file else OperationType.UPDATE
            message = f"成功{'创建' if is_new_file else '更新'}文件 '{req.path}'"
        
        # 创建文件信息
        file_info = fs_tool.create_file_info_from_path(
            p, 
            "文本文件" if is_new_file else "更新的文本文件",
            encoding=req.encoding,
            operation="append" if req.append else "write"
        )
        # 兼容序列化（支持 Pydantic/BaseModel/dataclass/dict）
        def _serialize_fi(fi):
            try:
                if fi is None:
                    return None
                if hasattr(fi, "model_dump") and callable(getattr(fi, "model_dump")):
                    return fi.model_dump()
                if hasattr(fi, "dict") and callable(getattr(fi, "dict")):
                    return fi.dict()
                if hasattr(fi, "to_dict") and callable(getattr(fi, "to_dict")):
                    return fi.to_dict()
                if isinstance(fi, dict):
                    return fi
            except Exception:
                pass
            # 最后兜底：返回字符串化路径信息
            try:
                return {"path": str(getattr(fi, "path", p)), "name": getattr(fi, "name", p.name)}
            except Exception:
                return {"path": str(p)}
        file_info_serialized = _serialize_fi(file_info)
        
        return builder.success(
            operation=operation_type,
            message=message,
            data={
                "path": req.path,
                "size": p.stat().st_size,
                "encoding": req.encoding,
                "is_new_file": is_new_file,
                "append_mode": req.append,
                "file_info": file_info_serialized
            },
            new_files=[file_info_serialized] if is_new_file else None
        ).to_dict()
        
    except Exception as e:
        return builder.error(
            operation=OperationType.WRITE,
            message="写入文件失败",
            error_details=str(e)
        ).to_dict()

@mcp.tool()
def read_binary(
    path: str = Field(..., description="要读取的二进制文件路径。适用于图片、音频、视频等非文本文件"),
    max_bytes: Optional[int] = Field(None, description="限制最大读取字节数。None表示读取整个文件，设置数值可避免读取过大文件")
) -> Dict[str, Any]:
    """读取二进制文件并返回base64编码的内容。
    
    用于读取图片、音频、视频、压缩包等二进制文件。文件内容会以base64格式返回，
    便于在JSON响应中传输和后续处理。
    
    Args:
        path: 要读取的二进制文件路径
        max_bytes: 可选的字节数限制，用于控制大文件的读取
    
    Returns:
        包含base64编码内容的响应，同时提供文件大小等元数据
    
    Examples:
        - read_binary("image.png") : 读取PNG图片文件
        - read_binary("audio.mp3", max_bytes=1024000) : 读取音频文件，限制1MB
        - read_binary("data.zip") : 读取压缩文件
    """
    builder = MCPResponseBuilder("read_binary")
    
    try:
        # 创建请求对象以保持兼容性
        req = ReadBinaryRequest(
            path=path,
            max_bytes=max_bytes
        )
        
        p = resolve_in_sandbox(req.path)
        if not p.exists() or not p.is_file():
            return builder.error(
                operation=OperationType.READ,
                message="文件不存在或不是文件",
                error_details=f"路径 '{req.path}' 不是一个有效文件"
            ).to_dict()
            
        data = p.read_bytes()
        if req.max_bytes is not None and req.max_bytes >= 0:
            data = data[:req.max_bytes]
            
        result = ReadBinaryResult(
            base64=base64.b64encode(data).decode("ascii"), 
            size=len(data)
        )
        
        return builder.success(
            operation=OperationType.READ,
            message=f"成功读取二进制文件 '{req.path}'",
            data=result.model_dump()
        ).to_dict()
        
    except Exception as e:
        return builder.error(
            operation=OperationType.READ,
            message="读取二进制文件失败",
            error_details=str(e)
        ).to_dict()

@mcp.tool()
def write_binary(
    path: str = Field(..., description="要创建的二进制文件路径。如果目录不存在会自动创建"),
    base64_data: str = Field(..., description="base64编码的二进制数据。可以是图片、音频、压缩包等文件的base64编码"),
    overwrite: bool = Field(True, description="如果文件已存在是否覆盖。True表示覆盖(默认)，False表示如果文件存在则报错")
) -> Dict[str, Any]:
    """从base64数据创建二进制文件。
    
    用于创建图片、音频、视频、压缩包等二进制文件。输入必须是有效的base64编码字符串。
    如果目标目录不存在，会自动创建所需的目录结构。
    
    Args:
        path: 要创建的二进制文件路径
        base64_data: 文件内容的base64编码字符串
        overwrite: 是否允许覆盖已存在的文件
    
    Returns:
        包含操作结果的响应，显示文件是否为新创建、文件大小等信息
    
    Examples:
        - write_binary("output.png", png_base64) : 创建PNG图片文件
        - write_binary("backup.zip", zip_base64, overwrite=False) : 创建压缩文件但不覆盖
        - write_binary("audio/sound.mp3", audio_base64) : 在子目录中创建音频文件
    """
    builder = MCPResponseBuilder("write_binary")
    
    try:
        # 创建请求对象以保持兼容性
        req = WriteBinaryRequest(
            path=path,
            base64=base64_data,
            overwrite=overwrite
        )
        
        p = resolve_in_sandbox(req.path)
        
        is_new_file = not p.exists()
        if p.exists() and not req.overwrite:
            return builder.error(
                operation=OperationType.WRITE,
                message="文件已存在",
                error_details=f"文件 '{req.path}' 已存在且未设置覆盖"
            ).to_dict()
        
        p.parent.mkdir(parents=True, exist_ok=True)
        data = base64.b64decode(req.base64)
        p.write_bytes(data)
        
        # 判断文件类型
        file_ext = p.suffix.lower()
        file_type_map = {
            '.jpg': 'JPEG图片', '.jpeg': 'JPEG图片', '.png': 'PNG图片', 
            '.gif': 'GIF图片', '.pdf': 'PDF文档', '.zip': '压缩文件',
            '.wav': '音频文件', '.mp3': 'MP3音频'
        }
        file_description = file_type_map.get(file_ext, '二进制文件')
        
        # 创建文件信息
        file_info = fs_tool.create_file_info_from_path(
            p, 
            file_description + ("" if is_new_file else " (已更新)"),
            operation_type="binary_write"
        )
        # 兼容序列化
        def _serialize_fi(fi):
            try:
                if fi is None:
                    return None
                if hasattr(fi, "model_dump") and callable(getattr(fi, "model_dump")):
                    return fi.model_dump()
                if hasattr(fi, "dict") and callable(getattr(fi, "dict")):
                    return fi.dict()
                if hasattr(fi, "to_dict") and callable(getattr(fi, "to_dict")):
                    return fi.to_dict()
                if isinstance(fi, dict):
                    return fi
            except Exception:
                pass
            try:
                return {"path": str(getattr(fi, "path", p)), "name": getattr(fi, "name", p.name)}
            except Exception:
                return {"path": str(p)}
        file_info_serialized = _serialize_fi(file_info)
        
        operation_type = OperationType.CREATE if is_new_file else OperationType.UPDATE
        message = f"成功{'创建' if is_new_file else '更新'}二进制文件 '{req.path}'"
        
        return builder.success(
            operation=operation_type,
            message=message,
            data={
                "path": req.path,
                "size": len(data),
                "is_new_file": is_new_file,
                "file_info": file_info_serialized
            },
            new_files=[file_info_serialized] if is_new_file else None
        ).to_dict()
        
    except Exception as e:
        return builder.error(
            operation=OperationType.WRITE,
            message="写入二进制文件失败",
            error_details=str(e)
        ).to_dict()

@mcp.tool()
def mkdir(
    path: str = Field(..., description="要创建的目录路径。可以是相对路径或绝对路径，支持创建多层嵌套目录"),
    parents: bool = Field(True, description="是否自动创建父目录。True表示会创建路径中所有不存在的父目录"),
    exist_ok: bool = Field(True, description="如果目录已存在的处理方式。True表示不报错(默认)，False表示已存在时抛出错误")
) -> Dict[str, Any]:
    """创建新目录。
    
    用于创建单个目录或多层嵌套目录结构。默认会自动创建所需的父目录，
    如果目录已存在也不会报错。
    
    Args:
        path: 要创建的目录路径
        parents: 是否自动创建缺失的父目录
        exist_ok: 目录已存在时的处理方式
    
    Returns:
        包含目录创建结果的响应，显示是否为新创建的目录
    
    Examples:
        - mkdir("new_folder") : 创建单个目录
        - mkdir("project/src/utils") : 创建多层嵌套目录
        - mkdir("temp", exist_ok=False) : 如果目录已存在则报错
        - mkdir("docs", parents=False) : 不自动创建父目录
    """
    builder = MCPResponseBuilder("mkdir")
    
    try:
        # 创建请求对象
        req = MkdirRequest(
            path=path,
            parents=parents,
            exist_ok=exist_ok
        )
        
        p = resolve_in_sandbox(req.path)
        is_new_dir = not p.exists()
        p.mkdir(parents=req.parents, exist_ok=req.exist_ok)
        
        if is_new_dir:
            # 创建目录信息
            file_info = fs_tool.create_file_info_from_path(
                p, 
                "目录",
                parents_created=req.parents
            )
            
            return builder.success(
                operation=OperationType.CREATE,
                message=f"成功创建目录 '{req.path}'",
                data={
                    "path": req.path,
                    "parents_created": req.parents
                },
                new_files=[file_info]
            ).to_dict()
        else:
            return builder.success(
                operation=OperationType.QUERY,
                message=f"目录 '{req.path}' 已存在",
                data={
                    "path": req.path,
                    "already_exists": True
                }
            ).to_dict()
            
    except FileExistsError:
        return builder.error(
            operation=OperationType.CREATE,
            message="目录已存在",
            error_details=f"目录 '{req.path}' 已存在且 exist_ok=False"
        ).to_dict()
    
    except Exception as e:
        return builder.error(
            operation=OperationType.CREATE,
            message="创建目录失败",
            error_details=str(e)
        ).to_dict()

@mcp.tool()
def move(
    src: str = Field(..., description="源文件或目录的路径。要移动的文件或目录必须存在"),
    dst: str = Field(..., description="目标路径。可以是新的文件名或目录路径"),
    overwrite: bool = Field(False, description="如果目标已存在是否覆盖。False表示不覆盖并报错(默认)，True表示覆盖现有文件")
) -> Dict[str, Any]:
    """移动文件或目录到新位置。
    
    将文件或目录从源位置移动到目标位置。这是一个剪切操作，源文件/目录在移动后将不再存在。
    支持重命名（在同一目录内移动）和跨目录移动。
    
    Args:
        src: 要移动的源文件或目录路径
        dst: 目标位置路径
        overwrite: 是否允许覆盖已存在的目标
    
    Returns:
        包含移动操作结果的响应，显示源路径、目标路径等信息
    
    Examples:
        - move("old_name.txt", "new_name.txt") : 重命名文件
        - move("temp.txt", "backup/temp.txt") : 移动文件到子目录
        - move("old_folder", "new_folder") : 重命名目录
        - move("file.txt", "existing.txt", overwrite=True) : 覆盖已存在的文件
    """
    builder = MCPResponseBuilder("move")
    
    try:
        # 创建请求对象以保持兼容性
        req = MoveCopyRequest(
            src=src,
            dst=dst,
            overwrite=overwrite
        )
        
        src_path = resolve_in_sandbox(req.src)
        dst_path = resolve_in_sandbox(req.dst)
        
        if not src_path.exists():
            return builder.error(
                operation=OperationType.UPDATE,
                message="源文件不存在",
                error_details=f"源路径 '{req.src}' 不存在"
            ).to_dict()
        
        # 检查目标是否已存在
        dst_existed = dst_path.exists()
        is_src_dir = src_path.is_dir()
        
        dst_path.parent.mkdir(parents=True, exist_ok=True)
        if dst_path.exists():
            if req.overwrite:
                if dst_path.is_dir():
                    shutil.rmtree(dst_path)
                else:
                    dst_path.unlink()
            else:
                return builder.error(
                    operation=OperationType.UPDATE,
                    message="目标已存在",
                    error_details=f"目标 '{req.dst}' 已存在，使用 overwrite=True 覆盖"
                ).to_dict()
        
        shutil.move(str(src_path), str(dst_path))
        
        # 创建文件信息
        file_info = fs_tool.create_file_info_from_path(
            dst_path,
            f"移动的{'目录' if is_src_dir else '文件'}",
            source=req.src,
            move_operation=True
        )
        
        return builder.success(
            operation=OperationType.UPDATE,
            message=f"成功移动{'目录' if is_src_dir else '文件'} '{req.src}' 到 '{req.dst}'",
            data={
                "source": req.src,
                "destination": req.dst,
                "type": "目录" if is_src_dir else "文件",
                "overwritten": dst_existed
            },
            new_files=[file_info] if not dst_existed else None,
            modified_files=[file_info] if dst_existed else None
        ).to_dict()
        
    except Exception as e:
        return builder.error(
            operation=OperationType.UPDATE,
            message="移动失败",
            error_details=str(e)
        ).to_dict()

@mcp.tool()
def copy(
    src: str = Field(..., description="源文件或目录的路径。要复制的文件或目录必须存在"),
    dst: str = Field(..., description="目标路径。复制后的新文件或目录位置"),
    overwrite: bool = Field(False, description="如果目标已存在是否覆盖。False表示不覆盖并报错(默认)，True表示覆盖现有文件")
) -> Dict[str, Any]:
    """复制文件或目录到新位置。
    
    创建文件或目录的副本。原始文件/目录保持不变，在目标位置创建相同的拷贝。
    支持文件复制和整个目录树的递归复制。
    
    Args:
        src: 要复制的源文件或目录路径
        dst: 目标位置路径
        overwrite: 是否允许覆盖已存在的目标
    
    Returns:
        包含复制操作结果的响应，显示源路径、目标路径等信息
    
    Examples:
        - copy("document.txt", "backup.txt") : 复制文件
        - copy("project", "project_backup") : 复制整个目录
        - copy("config.json", "configs/new_config.json") : 复制到子目录
        - copy("data.db", "old_data.db", overwrite=True) : 覆盖已存在的文件
    """
    builder = MCPResponseBuilder("copy")
    
    try:
        # 创建请求对象以保持兼容性
        req = MoveCopyRequest(
            src=src,
            dst=dst,
            overwrite=overwrite
        )
        
        src_path = resolve_in_sandbox(req.src)
        dst_path = resolve_in_sandbox(req.dst)
        
        if not src_path.exists():
            return builder.error(
                operation=OperationType.CREATE,
                message="源文件不存在",
                error_details=f"源路径 '{req.src}' 不存在"
            ).to_dict()
        
        # 检查目标是否已存在以及源的类型
        dst_existed = dst_path.exists()
        is_src_dir = src_path.is_dir()
        
        dst_path.parent.mkdir(parents=True, exist_ok=True)
        if dst_path.exists() and not req.overwrite:
            return builder.error(
                operation=OperationType.CREATE,
                message="目标已存在",
                error_details=f"目标 '{req.dst}' 已存在，使用 overwrite=True 覆盖"
            ).to_dict()
        
        if src_path.is_dir():
            if dst_path.exists():
                shutil.rmtree(dst_path)
            shutil.copytree(src_path, dst_path)
        else:
            shutil.copy2(src_path, dst_path)
        
        # 创建文件信息
        file_info = fs_tool.create_file_info_from_path(
            dst_path,
            f"复制的{'目录' if is_src_dir else '文件'}",
            source=req.src,
            copy_operation=True
        )
        
        return builder.success(
            operation=OperationType.CREATE,
            message=f"成功复制{'目录' if is_src_dir else '文件'} '{req.src}' 到 '{req.dst}'",
            data={
                "source": req.src,
                "destination": req.dst,
                "type": "目录" if is_src_dir else "文件",
                "overwritten": dst_existed
            },
            new_files=[file_info]
        ).to_dict()
        
    except Exception as e:
        return builder.error(
            operation=OperationType.CREATE,
            message="复制失败",
            error_details=str(e)
        ).to_dict()

@mcp.tool()
def delete(
    path: str = Field(..., description="要删除的文件或目录路径"),
    recursive: bool = Field(False, description="删除目录时是否递归删除。False表示只能删除空目录，True表示删除目录及其所有内容")
) -> Dict[str, Any]:
    """删除文件或目录。
    
    永久删除指定的文件或目录。对于目录，默认只能删除空目录，
    如需删除非空目录需要设置recursive=True。请谨慎使用此功能。
    
    Args:
        path: 要删除的文件或目录路径
        recursive: 是否递归删除非空目录
    
    Returns:
        包含删除操作结果的响应，显示删除的项目类型和大小等信息
    
    Examples:
        - delete("temp.txt") : 删除文件
        - delete("empty_folder") : 删除空目录
        - delete("project_folder", recursive=True) : 递归删除目录及其所有内容
    
    Warning:
        删除操作是不可逆的，请确认路径正确后再执行。
        recursive=True 会删除目录下的所有文件和子目录。
    """
    builder = MCPResponseBuilder("delete")
    
    try:
        # 创建请求对象以保持兼容性
        req = DeleteRequest(
            path=path,
            recursive=recursive
        )
        
        p = resolve_in_sandbox(req.path)
        
        if not p.exists():
            return builder.success(
                operation=OperationType.DELETE,
                message=f"目标 '{req.path}' 不存在，无需删除",
                data={
                    "path": req.path,
                    "existed": False
                }
            ).to_dict()
        
        # 记录删除的信息
        is_directory = p.is_dir()
        size = p.stat().st_size if p.is_file() else None
        
        # 执行删除
        if p.is_dir():
            if req.recursive:
                shutil.rmtree(p)
                item_type = "目录"
            else:
                p.rmdir()  # 只能删除空目录
                item_type = "空目录"
        else:
            p.unlink()
            item_type = "文件"
        
        return builder.success(
            operation=OperationType.DELETE,
            message=f"成功删除{item_type} '{req.path}'",
            data={
                "path": req.path,
                "type": item_type,
                "was_directory": is_directory,
                "size": size,
                "recursive": req.recursive
            },
            deleted_files=[fs_tool.get_relative_path(p)]
        ).to_dict()
        
    except OSError as e:
        return builder.error(
            operation=OperationType.DELETE,
            message="删除失败",
            error_details=f"无法删除 '{req.path}': {str(e)}"
        ).to_dict()
    
    except Exception as e:
        return builder.error(
            operation=OperationType.DELETE,
            message="删除操作失败",
            error_details=str(e)
        ).to_dict()

# ----------------------------- Office 文本提取 -------------------------------

def _read_pdf_text(p: Path) -> str:
    try:
        import pypdf  # type: ignore
    except Exception as e:  # pragma: no cover - 仅在未安装时触发
        raise RuntimeError("缺少依赖 pypdf，请使用: uv run ... --with pypdf") from e

    reader = pypdf.PdfReader(str(p))
    texts: list[str] = []
    for page in reader.pages:
        try:
            t = page.extract_text() or ""
        except Exception:
            t = ""
        if t:
            texts.append(t)
    return "\n\n".join(texts)


def _read_docx_text(p: Path) -> str:
    try:
        import docx  # type: ignore
    except Exception as e:  # pragma: no cover
        raise RuntimeError("缺少依赖 python-docx，请使用: uv run ... --with python-docx") from e

    d = docx.Document(str(p))
    parts: list[str] = []
    for para in d.paragraphs:
        if para.text:
            parts.append(para.text)
    return "\n".join(parts)


def _read_pptx_text(p: Path) -> str:
    try:
        from pptx import Presentation  # type: ignore
    except Exception as e:  # pragma: no cover
        raise RuntimeError("缺少依赖 python-pptx，请使用: uv run ... --with python-pptx") from e

    prs = Presentation(str(p))
    texts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            txt = None
            if hasattr(shape, "text"):
                txt = getattr(shape, "text")
            elif hasattr(shape, "has_text_frame") and getattr(shape, "has_text_frame"):
                tf = getattr(shape, "text_frame", None)
                if tf and hasattr(tf, "text"):
                    txt = tf.text
            if txt:
                texts.append(str(txt))
    return "\n\n".join(texts)


@mcp.tool()
def read_office_text(
    path: str = Field(..., description="Office文件路径，支持PDF(.pdf)、Word文档(.docx)、PowerPoint演示文稿(.pptx)格式"),
) -> Dict[str, Any]:
    """从Office文档中提取纯文本内容。
    
    从PDF、Word文档、PowerPoint演示文稿中提取纯文本内容，便于文档分析和处理。
    支持多页PDF、多段落Word文档和多张幻灯片的PowerPoint文件。
    
    Args:
        path: 要读取的Office文件路径，必须是.pdf、.docx或.pptx格式
    
    Returns:
        包含提取文本内容的响应，同时提供文档类型、文本长度等元数据
    
    Examples:
        - read_office_text("report.pdf") : 提取PDF文档文本
        - read_office_text("document.docx") : 提取Word文档文本  
        - read_office_text("presentation.pptx") : 提取PowerPoint文本
    
    Note:
        需要安装相应的依赖包：pypdf(PDF)、python-docx(Word)、python-pptx(PowerPoint)
    """
    builder = MCPResponseBuilder("read_office_text")
    
    try:
        # 创建请求对象以保持兼容性
        req = OfficeReadRequest(
            path=path
        )
        
        p = resolve_in_sandbox(req.path)
        if not p.exists() or not p.is_file():
            return builder.error(
                operation=OperationType.READ,
                message="文件不存在或不是文件",
                error_details=f"路径 '{req.path}' 不是一个有效文件"
            ).to_dict()
            
        ext = p.suffix.lower()
        if ext == ".pdf":
            text = _read_pdf_text(p)
            doc_type = "PDF文档"
        elif ext == ".docx":
            text = _read_docx_text(p)
            doc_type = "Word文档"
        elif ext == ".pptx":
            text = _read_pptx_text(p)
            doc_type = "PowerPoint文档"
        else:
            return builder.error(
                operation=OperationType.READ,
                message="不支持的文件格式",
                error_details="仅支持 .pdf/.docx/.pptx 文件"
            ).to_dict()
        
        return builder.success(
            operation=OperationType.PROCESS,
            message=f"成功提取{doc_type}文本内容",
            data={
                "path": req.path,
                "text": text,
                "document_type": doc_type,
                "text_length": len(text),
                "line_count": len(text.splitlines())
            }
        ).to_dict()
        
    except Exception as e:
        return builder.error(
            operation=OperationType.PROCESS,
            message="提取文本失败",
            error_details=str(e)
        ).to_dict()

# 保留同步功能但简化，因为它更多是系统功能而非核心MCP功能
# 这里只保留简化版本，重点在于展示如何将复杂功能适配到标准格式

# 简化的同步请求和结果模型
class SyncRequest(BaseModel):
    vm_id: str = Field(description="虚拟机ID")
    session_id: str = Field(description="会话ID") 
    target_base_path: str = Field(description="目标基础路径")

class SyncResult(BaseModel):
    success: bool
    message: str
    synced_files: List[str] = Field(default_factory=list)

@mcp.tool()
def sync_files_to_target(
    vm_id: str = Field(..., description="虚拟机标识符，用于区分不同的虚拟机实例"),
    session_id: str = Field(..., description="会话标识符，用于区分不同的用户会话"),
    target_base_path: str = Field(..., description="目标同步路径，文件将被同步到此路径下的子目录中")
) -> Dict[str, Any]:
    """将当前工作目录中的文件同步到指定的目标路径。
    
    将本地BASE_DIR中的文本文件（.txt, .md, .json, .yaml等）同步到目标路径。
    会在目标路径下创建以vm_id和session_id命名的子目录来存放同步的文件。
    
    Args:
        vm_id: 虚拟机标识符，用于创建唯一的同步目录
        session_id: 会话标识符，与vm_id一起确保目录唯一性
        target_base_path: 目标根路径，同步的文件将存放在此路径下
    
    Returns:
        包含同步操作结果的响应，显示同步的文件数量和文件列表
    
    Examples:
        - sync_files_to_target("vm1", "session1", "/backup") : 同步到/backup/vm1_session1/
        - sync_files_to_target("prod", "user123", "/shared/sync") : 同步到/shared/sync/prod_user123/
    
    Note:
        只同步文本格式文件，限制最多同步50个文件，会保持原有的目录结构
    """
    builder = MCPResponseBuilder("sync_files_to_target")
    
    try:
        # 创建请求对象以保持兼容性
        req = SyncRequest(
            vm_id=vm_id,
            session_id=session_id,
            target_base_path=target_base_path
        )
        
        # 这里是简化版本，实际可以根据需要扩展
        target_base_path_obj = Path(req.target_base_path) / f"{req.vm_id}_{req.session_id}"
        
        # 创建目标目录
        target_base_path_obj.mkdir(parents=True, exist_ok=True)
        
        # 扫描本地文件 - 简化版本只同步文本文件
        synced_files = []
        text_extensions = {'.txt', '.md', '.json', '.yaml', '.yml'}
        
        for file_path in BASE_DIR.rglob("*"):
            if file_path.is_file() and file_path.suffix.lower() in text_extensions:
                try:
                    relative_path = file_path.relative_to(BASE_DIR)
                    target_file_path = target_base_path_obj / relative_path
                    
                    # 确保目标父目录存在
                    target_file_path.parent.mkdir(parents=True, exist_ok=True)
                    
                    # 复制文件
                    shutil.copy2(file_path, target_file_path)
                    synced_files.append(str(relative_path))
                    
                    if len(synced_files) >= 50:  # 限制同步数量
                        break
                        
                except Exception:
                    continue
        
        result = SyncResult(
            success=True,
            message=f"成功同步 {len(synced_files)} 个文件到 {target_base_path_obj}",
            synced_files=synced_files
        )
        
        return builder.success(
            operation=OperationType.SYSTEM,
            message=result.message,
            data=result.model_dump()
        ).to_dict()
        
    except Exception as e:
        return builder.error(
            operation=OperationType.SYSTEM,
            message="同步失败",
            error_details=str(e)
        ).to_dict()

# -----------------------------------------------------------------------------
# 启动
# -----------------------------------------------------------------------------

if __name__ == "__main__":
    # 保持兼容性，使用原有启动方式
    import sys
    import os
    sys.path.append(os.path.join(os.path.dirname(__file__), '../..'))
    from server_base import start_mcp_server
    
    # 添加自定义HTTP端点 - 正确的异步版本
    @mcp.custom_route("/direct/list-all-paths", ["GET"])
    async def http_list_all_paths(request):
        """HTTP端点：获取所有路径列表 - 专用于直接调用"""
        from fastapi.responses import JSONResponse
        
        try:
            # 调用完整的list_all_paths函数
            result = list_all_paths()
            
            # 简化返回格式，只返回必要的字段以确保兼容性
            if isinstance(result, dict) and result.get('status') == 'success':
                response_data = {
                    "success": True,
                    "data": result.get('data', {}),
                    "message": result.get('message', '')
                }
            else:
                response_data = {
                    "success": False,
                    "data": {},
                    "message": "获取路径列表失败"
                }
                
            return JSONResponse(content=response_data)
            
        except Exception as e:
            import traceback
            traceback.print_exc()
            # 返回JSON格式的错误
            error_response = {
                "success": False,
                "data": {},
                "message": f"获取路径列表失败: {str(e)}"
            }
            return JSONResponse(content=error_response, status_code=500)
    
    print(f"🚀 启动标准化文件系统服务器")
    print(f"📁 基础目录: {BASE_DIR}")
    print(f"🌐 端口: 8003")
    print(f"📋 使用标准响应格式")
    print(f"🔗 HTTP端点: /direct/list-all-paths (专用于直接调用)")
    
    start_mcp_server(mcp, 8003, "filesystem")